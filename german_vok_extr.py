#!/usr/bin/env python3
"""
German Vocabulary Extractor - Lite Version (No spaCy)
Handles PDF/DOCX/PPTX with OCR, basic analysis, deduplication & review workflow
"""

import os
import re
import sys
import csv
import json
from pathlib import Path
from collections import defaultdict
from typing import List, Dict, Set, Tuple
import argparse
import logging
from tqdm import tqdm

# External dependencies
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import pandas as pd
from langdetect import detect, DetectorFactory

# Ensure consistent results from langdetect
DetectorFactory.seed = 0

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s | %(levelname)-8s | %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

class GermanVocabExtractor:
    """Lightweight German vocabulary extraction engine (no spaCy)"""
    
    # Common German articles and prepositions
    ARTICLES = {'der', 'die', 'das', 'den', 'dem', 'des', 'ein', 'eine', 'einen', 'einem', 'eines'}
    CASE_PREPOSITIONS = {
        'akkusativ': {'für', 'durch', 'gegen', 'ohne', 'um', 'bis'},
        'dativ': {'mit', 'nach', 'von', 'zu', 'aus', 'bei', 'seit', 'gegenüber'},
        'genitiv': {'wegen', 'trotz', 'während', 'anstatt', 'innerhalb', 'außerhalb'}
    }
    # Basic words to exclude
    STOP_WORDS = {
        'der', 'die', 'das', 'ein', 'eine', 'und', 'oder', 'aber', 'denn', 'sondern',
        'als', 'wie', 'wenn', 'weil', 'dass', 'zu', 'in', 'an', 'auf', 'unter', 'über',
        'nein', 'ja', 'bitte', 'danke', 'hallo', 'ich', 'du', 'er', 'sie', 'es', 'wir',
        'ihr', 'mein', 'dein', 'sein', 'unser', 'euer', 'ist', 'hat', 'war', 'sind',
        'haben', 'hatte', 'wird', 'wurde', 'kann', 'konnte', 'muss', 'musste', 'soll'
    }
    
    def __init__(self, min_word_length: int = 3, force_german: bool = False):
        self.min_word_length = min_word_length
        self.force_german = force_german
        self.vocab_set: Set[str] = set()
        self.vocab_data: Dict[str, Dict] = {}
        self.uncertain_words: List[Dict] = []
        self.expressions: Set[str] = set()
        self.processed_files: List[str] = []
        self.stats = defaultdict(int)
    
    def extract_from_pdf(self, filepath: str) -> str:
        """Extract text from PDF with OCR fallback for images"""
        text = ""
        try:
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                        self.stats['pdf_pages_text'] += 1
                    
                    images = page.images
                    if images:
                        self.stats['pdf_pages_with_images'] += 1
                        for img in images:
                            try:
                                im = page.crop((img['x0'], img['top'], img['x1'], img['bottom'])).to_image()
                                img_pil = im.original
                                ocr_text = pytesseract.image_to_string(
                                    img_pil, lang='deu', config='--psm 6'
                                )
                                if ocr_text.strip():
                                    text += ocr_text + "\n"
                                    self.stats['ocr_images_processed'] += 1
                            except Exception as e:
                                logger.warning(f"OCR failed on image in {filepath} page {page_num+1}: {e}")
                                self.stats['ocr_failures'] += 1
            
            if not text.strip():
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        try:
                            text += page.extract_text() + "\n"
                            self.stats['pdf_fallback_used'] += 1
                        except Exception as e:
                            logger.warning(f"PyPDF2 extraction failed: {e}")
            
            self.stats['pdf_files_processed'] += 1
            logger.info(f"Extracted {len(text)} chars from PDF: {filepath}")
        except Exception as e:
            logger.error(f"PDF processing failed for {filepath}: {e}")
            self.stats['pdf_processing_errors'] += 1
        return text
    
    def extract_from_docx(self, filepath: str) -> str:
        """Extract text from DOCX files"""
        text = ""
        try:
            doc = Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\n"
            self.stats['docx_files_processed'] += 1
            logger.info(f"Extracted {len(text)} chars from DOCX: {filepath}")
        except Exception as e:
            logger.error(f"DOCX processing failed for {filepath}: {e}")
            self.stats['docx_processing_errors'] += 1
        return text
    
    def extract_from_pptx(self, filepath: str) -> str:
        """Extract text from PPTX files"""
        text = ""
        try:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            self.stats['pptx_files_processed'] += 1
            logger.info(f"Extracted {len(text)} chars from PPTX: {filepath}")
        except Exception as e:
            logger.error(f"PPTX processing failed for {filepath}: {e}")
            self.stats['pptx_processing_errors'] += 1
        return text
    
    def tokenize_text(self, text: str) -> List[str]:
        """Simple tokenization using regex (no spaCy)"""
        tokens = re.findall(r'\b[a-zA-ZäöüßÄÖÜ][a-zA-ZäöüßÄÖÜ\-]*\b', text)
        return tokens
    
    def normalize_word(self, word: str) -> str:
        """Normalize word (lowercase, clean punctuation)"""
        word = word.strip().lower()
        word = re.sub(r'[^\wäöüß\-]', '', word)
        if not word or len(word) < self.min_word_length:
            return ""
        return word
    
    def is_german_text(self, text: str, min_length: int = 50) -> bool:
        """Verify text is German using language detection"""
        if self.force_german:
            return True
        
        if len(text) < min_length:
            return True
        
        chunks = []
        text_len = len(text)
        chunks.append(text[:500])
        
        if text_len > 1000:
            chunks.append(text[text_len//2 : text_len//2 + 500])
            chunks.append(text[-500:])
        
        german_count = 0
        total_checks = 0
        
        for chunk in chunks:
            if len(chunk) < 50:
                continue
            try:
                lang = detect(chunk)
                total_checks += 1
                if lang == 'de':
                    german_count += 1
            except:
                continue
        
        if total_checks == 0:
            return True
        
        return (german_count / total_checks) >= 0.5
    
    def analyze_word_simple(self, word: str, context: str = "") -> Dict:
        """Simple word analysis without spaCy"""
        analysis = {
            'word': word,
            'lemma': self.normalize_word(word),
            'article': '',
            'plural': '',
            'case_usage': [],
            'pos': 'UNKNOWN',
            'confidence': 0.7
        }
        
        # Detect article from context
        if context:
            words = context.lower().split()
            try:
                idx = words.index(word.lower())
                if idx > 0:
                    prev_word = words[idx-1]
                    if prev_word in self.ARTICLES:
                        analysis['article'] = prev_word
                        analysis['pos'] = 'NOUN'
                        analysis['confidence'] = 0.9
            except ValueError:
                pass
        
        # Simple plural heuristic
        if analysis['pos'] == 'NOUN' or word[0].isupper():
            if word.endswith('e'):
                analysis['plural'] = word + 'n'
            elif word.endswith('el') or word.endswith('er'):
                analysis['plural'] = word + 'n'
            else:
                analysis['plural'] = word + 'e'
        
        # Case usage from context
        if context:
            context_lower = context.lower()
            for case, preps in self.CASE_PREPOSITIONS.items():
                if any(prep in context_lower for prep in preps):
                    analysis['case_usage'].append(case)
        
        return analysis
    
    def detect_expressions_simple(self, text: str) -> Set[str]:
        """Detect common German expressions via regex"""
        expressions = set()
        
        expression_patterns = [
            r'\b(wie geht es)\b',
            r'\b(es tut mir leid)\b',
            r'\b(vielen dank)\b',
            r'\b(bitte schön)\b',
            r'\b(kein problem)\b',
            r'\b(bis bald)\b',
            r'\b(auf wiedersehen)\b',
            r'\b(guten tag)\b',
            r'\b(gute nacht)\b',
            r'\b(guten morgen)\b',
            r'\b(guten abend)\b',
        ]
        
        for pattern in expression_patterns:
            matches = re.findall(pattern, text.lower())
            expressions.update(matches)
        
        return expressions
    
    def should_include_word(self, word: str, analysis: Dict) -> Tuple[bool, str]:
        """Determine if word should be included"""
        if len(word) < self.min_word_length:
            return False, "too_short"
        
        if word.lower() in self.STOP_WORDS:
            return False, "stop_word"
        
        if not re.match(r'^[a-zA-ZäöüßÄÖÜ\-]+$', word):
            return False, "non_german_chars"
        
        if word.isdigit() or word.replace(',', '').replace('.', '').isdigit():
            return False, "number"
        
        # Capitalized words might be proper nouns (uncertain)
        if word[0].isupper() and analysis['article'] == '':
            return False, "possible_proper_noun"
        
        return True, "accepted"
    
    def process_text(self, text: str, source: str):
        """Process extracted text and populate vocabulary sets"""
        if not text.strip():
            logger.warning(f"No text extracted from {source}")
            return
        
        if not self.is_german_text(text):
            logger.warning(f"Skipping non-German content in {source}")
            self.stats['non_german_content_skipped'] += 1
            return
        
        # Detect expressions
        expressions = self.detect_expressions_simple(text)
        for expr in expressions:
            norm_expr = expr.strip().lower()
            if norm_expr not in self.expressions:
                self.expressions.add(norm_expr)
                self.stats['expressions_found'] += 1
        
        # Tokenize and analyze
        tokens = self.tokenize_text(text)
        tokens_processed = 0
        
        for word in tqdm(tokens, desc=f"Processing {Path(source).name}", leave=False):
            if not word or len(word) < self.min_word_length:
                continue
            
            # Get context (surrounding words)
            try:
                idx = tokens.index(word)
            except ValueError:
                idx = 0
            start = max(0, idx - 3)
            end = min(len(tokens), idx + 4)
            context = ' '.join(tokens[start:end])
            
            analysis = self.analyze_word_simple(word, context)
            norm_word = analysis['lemma']
            
            if not norm_word:
                continue
            
            tokens_processed += 1
            
            if norm_word in self.vocab_set:
                self.stats['duplicates_skipped'] += 1
                continue
            
            include, reason = self.should_include_word(word, analysis)
            
            if include:
                self.vocab_set.add(norm_word)
                self.vocab_data[norm_word] = {
                    'original': word,
                    'lemma': norm_word,
                    'article': analysis['article'],
                    'plural': analysis['plural'],
                    'case_usage': ', '.join(analysis['case_usage']) if analysis['case_usage'] else '',
                    'pos': analysis['pos'],
                    'source': source,
                    'context': context[:100] + '...' if len(context) > 100 else context
                }
                self.stats['words_accepted'] += 1
            else:
                self.uncertain_words.append({
                    'word': word,
                    'reason': reason,
                    'lemma': norm_word,
                    'pos': analysis['pos'],
                    'source': source,
                    'context': context[:80] + '...'
                })
                self.stats[f'uncertain_{reason}'] += 1
        
        self.stats['tokens_processed'] += tokens_processed
        self.processed_files.append(source)
        logger.info(f"Processed {tokens_processed} tokens from {source}")
    
    def show_review_table(self) -> bool:
        """Display summary table and prompt for user confirmation"""
        print("\n" + "="*80)
        print("🇩🇪 GERMAN VOCABULARY EXTRACTION SUMMARY")
        print("="*80)
        
        if self.vocab_data:
            df = pd.DataFrame(self.vocab_data.values())
            display_df = df[['original', 'article', 'plural', 'case_usage', 'pos']].head(20)
            print("\n✅ ACCEPTED VOCABULARY (First 20 entries):")
            print(display_df.to_string(index=False))
            print(f"\nTotal unique words: {len(self.vocab_data)}")
        else:
            print("\n⚠️  No vocabulary extracted!")
        
        if self.uncertain_words:
            uncertain_df = pd.DataFrame(self.uncertain_words[:15])
            print("\n❓ UNCERTAIN WORDS (Require review):")
            print(uncertain_df[['word', 'reason', 'pos', 'context']].to_string(index=False))
            print(f"\nTotal uncertain words: {len(self.uncertain_words)}")
        
        print("\n📊 PROCESSING STATISTICS:")
        for key in sorted(self.stats.keys()):
            print(f"  • {key.replace('_', ' ').title()}: {self.stats[key]}")
        
        print("\n" + "="*80)
        print("ACTIONS:")
        print("  [C]ontinue - Save results and exit")
        print("  [R]etry    - Restart extraction (preserves current vocabulary)")
        print("  [Q]uit     - Abort without saving")
        print("="*80)
        
        while True:
            choice = input("\nYour choice [C/R/Q]: ").strip().lower()
            if choice == 'c':
                return True
            elif choice == 'r':
                return False
            elif choice == 'q':
                sys.exit(0)
            else:
                print("Invalid choice. Please enter C, R, or Q.")
    
    def save_results(self, output_dir: str = "output"):
        """Save vocabulary and uncertain words to TXT and CSV formats"""
        Path(output_dir).mkdir(exist_ok=True)
        
        timestamp = pd.Timestamp.now().strftime("%Y%m%d_%H%M%S")
        vocab_items = sorted(self.vocab_data.values(), key=lambda x: x['lemma'])
        
        # TXT format
        txt_path = Path(output_dir) / f"vocabulary_{timestamp}.txt"
        with open(txt_path, 'w', encoding='utf-8') as f:
            for item in vocab_items:
                line_parts = [item['original']]
                if item['article']:
                    line_parts.append(f"({item['article']})")
                if item['plural']:
                    line_parts.append(f"Pl: {item['plural']}")
                if item['case_usage']:
                    line_parts.append(f"[{item['case_usage']}]")
                f.write(' '.join(line_parts) + '\n')
        
        # CSV format - FIXED: Added 'context' to fieldnames
        csv_path = Path(output_dir) / f"vocabulary_{timestamp}.csv"
        with open(csv_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=[
                'original', 'lemma', 'article', 'plural', 'case_usage', 'pos', 'source', 'context'
            ])
            writer.writeheader()
            writer.writerows(vocab_items)
        
        # Uncertain words
        if self.uncertain_words:
            uncertain_txt = Path(output_dir) / f"uncertain_words_{timestamp}.txt"
            with open(uncertain_txt, 'w', encoding='utf-8') as f:
                for item in self.uncertain_words:
                    f.write(f"{item['word']} | Reason: {item['reason']} | Context: {item['context']}\n")
            
            uncertain_csv = Path(output_dir) / f"uncertain_words_{timestamp}.csv"
            with open(uncertain_csv, 'w', encoding='utf-8', newline='') as f:
                writer = csv.DictWriter(f, fieldnames=['word', 'reason', 'lemma', 'pos', 'source', 'context'])
                writer.writeheader()
                writer.writerows(self.uncertain_words)
        
        # Expressions
        if self.expressions:
            expr_path = Path(output_dir) / f"expressions_{timestamp}.txt"
            with open(expr_path, 'w', encoding='utf-8') as f:
                for expr in sorted(self.expressions):
                    f.write(expr + '\n')
        
        # Report
        report_path = Path(output_dir) / f"extraction_report_{timestamp}.json"
        report = {
            'timestamp': timestamp,
            'files_processed': self.processed_files,
            'statistics': dict(self.stats),
            'vocabulary_count': len(self.vocab_data),
            'uncertain_count': len(self.uncertain_words),
            'expressions_count': len(self.expressions)
        }
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        
        logger.info(f"Results saved to {output_dir}/")
        print(f"\n✅ Saved {len(vocab_items)} unique words to:")
        print(f"   • TXT: {txt_path.name}")
        print(f"   • CSV: {csv_path.name}")
        if self.uncertain_words:
            print(f"   • Uncertain words: {len(self.uncertain_words)} entries")
        if self.expressions:
            print(f"   • Expressions: {len(self.expressions)} phrases")
    
    def process_files(self, filepaths: List[str]):
        """Main processing pipeline"""
        for filepath in filepaths:
            filepath = Path(filepath)
            if not filepath.exists():
                logger.error(f"File not found: {filepath}")
                self.stats['missing_files'] += 1
                continue
            
            logger.info(f"Processing: {filepath.name}")
            text = ""
            
            try:
                if filepath.suffix.lower() == '.pdf':
                    text = self.extract_from_pdf(str(filepath))
                elif filepath.suffix.lower() == '.docx':
                    text = self.extract_from_docx(str(filepath))
                elif filepath.suffix.lower() == '.pptx':
                    text = self.extract_from_pptx(str(filepath))
                else:
                    logger.warning(f"Unsupported file type: {filepath.suffix}")
                    self.stats['unsupported_files'] += 1
                    continue
                
                if text.strip():
                    self.process_text(text, str(filepath))
                else:
                    logger.warning(f"No extractable text from {filepath.name}")
                    self.stats['empty_extractions'] += 1
            
            except Exception as e:
                logger.error(f"Processing failed for {filepath}: {e}")
                self.stats['processing_exceptions'] += 1
        
        while True:
            proceed = self.show_review_table()
            if proceed:
                break
            else:
                self.stats.clear()
                logger.info("Restarting extraction (existing vocabulary preserved)")
        
        self.save_results()

def main():
    parser = argparse.ArgumentParser(
        description="🇩🇪 German Vocabulary Extractor (Lite - No spaCy)",
        formatter_class=argparse.RawTextHelpFormatter
    )
    parser.add_argument('files', nargs='+', help='Input files (PDF, DOCX, PPTX)')
    parser.add_argument('--min-length', type=int, default=3, 
                       help='Minimum word length (default: 3)')
    parser.add_argument('--output-dir', default='output', 
                       help='Output directory (default: output)')
    parser.add_argument('--no-ocr', action='store_true', 
                       help='Disable OCR')
    parser.add_argument('--force', action='store_true',
                       help='Force processing (bypass language check)')
    parser.add_argument('--verbose', action='store_true', 
                       help='Enable debug logging')
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    if args.no_ocr:
        logger.warning("OCR disabled")
    
    extractor = GermanVocabExtractor(min_word_length=args.min_length, force_german=args.force)
    extractor.process_files(args.files)
    
    print("\n✨ Extraction completed successfully!")
    print(f"   Vocabulary files saved in: {args.output_dir}/")

if __name__ == "__main__":
    main()